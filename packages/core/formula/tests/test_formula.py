"""Unit tests for formula rendering.

Each test maps to a scenario in
openspec/changes/add-formula-rendering/specs/formula-rendering/spec.md.
"""

from __future__ import annotations

import pytest

from formula_render import MatplotlibFormulaRenderer

_PNG_SIG = b"\x89PNG\r\n\x1a\n"


def _renderer(tmp_path):
    # Lower dpi keeps tests fast; behavior is identical.
    return MatplotlibFormulaRenderer(out_dir=tmp_path, dpi=150)


def test_parseable_returns_png(tmp_path):
    p = _renderer(tmp_path).to_image("E=mc^2")
    assert p is not None and p.exists()
    assert p.read_bytes()[:8] == _PNG_SIG


def test_unparseable_returns_none(tmp_path):
    # mhchem is not supported by mathtext -> graceful None
    assert _renderer(tmp_path).to_image(r"\ce{H2O}") is None


@pytest.mark.parametrize(
    "latex",
    [
        "^{143}Nd/^{144}Nd",
        "SiO_2",
        r"\frac{1}{n}\sum_{i=1}^{n}(Y_i-\hat{Y}_i)^2",
        r"\varepsilon_{Nd}(t)",
        r"\sqrt{MSE}",
    ],
)
def test_academic_formulas_render(tmp_path, latex):
    p = _renderer(tmp_path).to_image(latex)
    assert p is not None and p.exists()


def test_caching_reuses_image(tmp_path):
    r = _renderer(tmp_path)
    p1 = r.to_image("a^2+b^2")
    assert p1 is not None
    mtime1 = p1.stat().st_mtime_ns
    p2 = r.to_image("a^2+b^2")
    assert p2 == p1
    assert p2.stat().st_mtime_ns == mtime1  # not re-rendered


def test_compiler_embeds_picture(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx_compiler import compile_deck
    from slide_ir import Deck, FormulaBlock, LayoutType, SlideIR

    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FORMULA_BANNER,
                title="f",
                blocks=[FormulaBlock(latex="E=mc^2")],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "o.pptx", formula_renderer=_renderer(tmp_path))
    slide = Presentation(str(out)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


# --- enhance-batch3-formula: MathJax tier ------------------------------------


def test_is_advanced_detection():
    from formula_render import is_advanced

    assert is_advanced(r"\ce{2H2 + O2 -> 2H2O}")
    assert is_advanced(r"\begin{pmatrix} a & b \ c & d \end{pmatrix}")
    assert not is_advanced("x^2 + y^2")
    assert not is_advanced(r"\frac{a}{b}")


def test_auto_routes_advanced_to_mathjax(tmp_path):
    from formula_render import AutoFormulaRenderer

    calls = {"adv": 0, "simple": 0}

    class _Adv:
        def to_image(self, latex):
            calls["adv"] += 1
            return tmp_path / "a.png"

    class _Simple:
        def to_image(self, latex):
            calls["simple"] += 1
            return tmp_path / "s.png"

    (tmp_path / "a.png").write_bytes(b"x")
    (tmp_path / "s.png").write_bytes(b"x")
    r = AutoFormulaRenderer(simple=_Simple(), advanced=_Adv())
    r.to_image(r"\ce{H2O}")  # advanced -> mathjax
    r.to_image("x^2")  # simple -> matplotlib
    assert calls["adv"] == 1 and calls["simple"] == 1


def test_mathjax_renders_chemistry_if_available(tmp_path):
    import pytest

    from formula_render import MathJaxFormulaRenderer

    if not MathJaxFormulaRenderer.available():
        pytest.skip("Node formula sidecar not installed (npm install in packages/core/formula/node)")
    out = MathJaxFormulaRenderer(tmp_path).to_image(r"\ce{2H2 + O2 -> 2H2O}")
    assert out and out.exists() and out.stat().st_size > 0
