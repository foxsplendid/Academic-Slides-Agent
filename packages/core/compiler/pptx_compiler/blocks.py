"""Render individual Slide-IR blocks into native python-pptx shapes within a region.

A region is an EMU tuple ``(left, top, width, height)``. Text is rendered CJK-aware (East-Asian +
Latin typefaces) and supports a ``**…**`` emphasis convention rendered as bold red runs, matching the
reference 组会 deck style (docs/SPEC.md styling).
"""

from __future__ import annotations

import math
import re
from pathlib import Path

from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from slide_ir import BulletBlock, ChartBlock, FigureBlock, FormulaBlock, TableBlock

from .formula_renderer import FormulaRenderer
from .style import ACADEMIC, StyleProfile

Region = tuple[int, int, int, int]

EA_FONT = ACADEMIC.ea_font  # kept for back-compat; per-render styling comes from the StyleProfile
LATIN_FONT = ACADEMIC.latin_font
_EMPHASIS = re.compile(r"\*\*(.+?)\*\*")


def _set_ea(font, typeface: str) -> None:
    """Set the East-Asian typeface (`<a:ea>`); python-pptx's ``font.name`` only sets `<a:latin>`."""
    rpr = font._rPr
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", typeface)


def add_rich_text(paragraph, text: str, *, size, bold: bool = False, style: StyleProfile = ACADEMIC, color=None) -> None:
    """Add runs to ``paragraph``, rendering ``**…**`` spans as bold emphasis-colored while the rest uses
    ``color`` (or the theme default when None)."""
    segments: list[tuple[str, bool]] = []
    pos = 0
    for m in _EMPHASIS.finditer(text):
        if m.start() > pos:
            segments.append((text[pos:m.start()], False))
        segments.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        segments.append((text[pos:], False))
    if not segments:
        segments = [(text, False)]
    for content, emph in segments:
        run = paragraph.add_run()
        run.text = content
        f = run.font
        f.size = size
        f.bold = bold or emph
        f.name = style.latin_font
        if emph:
            f.color.rgb = style.emphasis_rgb
        elif color is not None:
            f.color.rgb = color
        _set_ea(f, style.ea_font)


_EMU_PER_PT = 12700.0


def _display_width(text: str) -> float:
    """CJK chars occupy ~1 em, Latin ~0.5 em — used to estimate wrapped line count."""
    return sum(1.0 if ord(c) > 0x2E80 else 0.5 for c in text)


def _est_lines(items: list[str], w_pt: float, f: float) -> int:
    per_line = max(w_pt / f, 1.0)  # ~em units that fit on one line at size f
    return sum(max(1, math.ceil(_display_width("• " + it) / per_line)) for it in items)


def _fit_font(items: list[str], width: int, height: int, *, base: float = 16.0, floor: float = 10.0, grow: float = 0.0) -> float:
    """Measure, then place: shrink the font until the bullets' estimated height fits the region —
    or, when ``grow`` allows and the text fills under ~55% of the region, bump it up so short lists
    don't leave the bottom half of a page empty (the R6' whitespace complaint)."""
    w_pt = max(width / _EMU_PER_PT, 1.0)
    h_pt = max(height / _EMU_PER_PT, 1.0)
    f = base
    while f > floor:
        if _est_lines(items, w_pt, f) * f * 1.28 <= h_pt:
            break
        f -= 1.0
    while grow > 0 and f < base + grow:
        nxt = f + 1.0
        if _est_lines(items, w_pt, nxt) * nxt * 1.28 <= h_pt * 0.55:
            f = nxt
        else:
            break
    return max(f, floor)


def _flat_bullets(items) -> list[tuple[str, int]]:
    """Flatten BulletBlock items (str | BulletItem) into (text, level) pairs."""
    out: list[tuple[str, int]] = []
    for it in items:
        if isinstance(it, str):
            out.append((it, 0))
        else:  # BulletItem
            out.append((it.text, 0))
            out.extend((c, 1) for c in it.children)
    return out


def _set_bullet_props(para, level: int, size_pt: float) -> None:
    """Real PPT bullet formatting: hanging indent + bullet glyph per level (•, then –)."""
    pPr = para._p.get_or_add_pPr()
    hang = int(Pt(size_pt) * 1.15)
    pPr.set("marL", str(hang * (level + 1)))
    pPr.set("indent", str(-hang))
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "•" if level == 0 else "–"})
    pPr.append(buFont)
    pPr.append(buChar)


def render_bullets(slide, block: BulletBlock, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    flat = _flat_bullets(block.items)
    # 8% bottom safety margin: the estimator can undershoot on mixed CJK/Latin runs, and an
    # overflowing textbox spills into the region below it (seen as text clipped by a callout).
    size_pt = _fit_font([t for t, _ in flat], width, int(height * 0.92), base=style.body_pt, grow=6.0)
    for i, (text, level) in enumerate(flat):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        sub = level > 0
        add_rich_text(p, text, size=Pt(size_pt - 2 if sub else size_pt), style=style)
        _set_bullet_props(p, level, size_pt)
    return box


def _icon_png(icon_resolver, name, style: StyleProfile, px: int = 64):
    """Resolve an icon name to a tinted PNG path via the injected resolver; None = skip (fail open)."""
    if not (icon_resolver and name):
        return None
    try:
        return icon_resolver(name, f"#{style.accent_rgb}", px)
    except Exception:
        return None


def render_callout(slide, block, region: Region, style: StyleProfile = ACADEMIC, icon_resolver=None):
    """A tinted takeaway card with an accent edge on the left (and an optional concept icon)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    left, top, width, height = region
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = style.card_fill_rgb
    card.line.color.rgb = style.card_line_rgb
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, int(Pt(5)), height)
    edge.fill.solid()
    edge.fill.fore_color.rgb = style.accent_rgb
    edge.line.fill.background()
    edge.shadow.inherit = False

    icon = _icon_png(icon_resolver, getattr(block, "icon", None), style)
    if icon is not None:
        size = int(Pt(20))
        slide.shapes.add_picture(str(icon), left + int(Pt(12)), top + (height - size) // 2, size, size)

    tf = card.text_frame
    if icon is not None:
        tf.margin_left = int(Pt(40))  # clear the icon
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    if block.label:
        tag = para.add_run()
        tag.text = f"{block.label}  "
        tag.font.bold = True
        tag.font.size = Pt(style.body_pt)
        tag.font.name = style.latin_font
        tag.font.color.rgb = style.accent_rgb
        _set_ea(tag.font, style.ea_font)
    # Explicit body color: shape-style text defaults to white and would vanish on the light card.
    add_rich_text(para, block.text, size=Pt(style.body_pt), style=style, color=style.text_rgb)
    return card


def render_stat(slide, block, region: Region, style: StyleProfile = ACADEMIC, icon_resolver=None):
    """1-4 big-number cards in a row (value in accent color, label muted below, optional icon)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR

    left, top, width, height = region
    n = len(block.items)
    gap = int(Pt(10))
    card_w = (width - gap * (n - 1)) // n
    value_pt = min(34.0, max(22.0, height / Pt(1) * 0.32))  # scale with the band height
    for i, item in enumerate(block.items):
        x = left + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, height)
        card.fill.solid()
        card.fill.fore_color.rgb = style.card_fill_rgb
        card.line.color.rgb = style.card_line_rgb
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + int(Pt(6)), top, card_w - int(Pt(12)), int(Pt(3)))
        edge.fill.solid()
        edge.fill.fore_color.rgb = style.accent_rgb
        edge.line.fill.background()
        edge.shadow.inherit = False
        icon = _icon_png(icon_resolver, getattr(item, "icon", None), style)
        if icon is not None:
            size = int(Pt(16))
            slide.shapes.add_picture(str(icon), x + int(Pt(8)), top + int(Pt(8)), size, size)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = tf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        run = vp.add_run()
        run.text = item.value
        run.font.bold = True
        run.font.size = Pt(value_pt)
        run.font.name = style.latin_font
        run.font.color.rgb = style.accent_rgb
        _set_ea(run.font, style.ea_font)
        if item.label:
            lp = tf.add_paragraph()
            lp.alignment = PP_ALIGN.CENTER
            lr = lp.add_run()
            lr.text = item.label
            lr.font.size = Pt(style.caption_pt)
            lr.font.name = style.latin_font
            lr.font.color.rgb = style.muted_rgb
            _set_ea(lr.font, style.ea_font)
    return None


def render_table(slide, block: TableBlock, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n_cols = len(block.columns)
    n_rows = len(block.rows) + 1  # header + data
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = graphic_frame.table
    # `highlight` may mark cells {"cells": [[row, col], ...]} (data-row coordinates, 0-based).
    marked = set()
    if isinstance(block.highlight, dict):
        for rc in block.highlight.get("cells", []):
            if isinstance(rc, (list, tuple)) and len(rc) == 2:
                marked.add((int(rc[0]), int(rc[1])))

    white = RGBColor(0xFF, 0xFF, 0xFF)
    for c, name in enumerate(block.columns):
        cell = table.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = style.table_header_rgb
        add_rich_text(
            cell.text_frame.paragraphs[0], str(name), size=Pt(style.table_header_pt), bold=True, style=style, color=white
        )

    for r, row in enumerate(block.rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.text = ""
            if r % 2 == 0:  # zebra banding on even data rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = style.table_band_rgb
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = white
            value = str(row[c]) if c < len(row) else ""
            hot = (r - 1, c) in marked
            add_rich_text(
                cell.text_frame.paragraphs[0],
                value,
                size=Pt(style.table_body_pt),
                bold=hot,
                style=style,
                color=style.emphasis_rgb if hot else None,
            )

    return graphic_frame


_CATEGORY_CHARTS = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


def _style_chart(chart, block: ChartBlock, style: StyleProfile) -> None:
    """Apply StyleProfile design tokens to a native chart (palette, fonts, labels, legend) so it
    matches the deck instead of looking like a raw Office default. Best-effort: any sub-step that a
    given chart type doesn't support is skipped."""
    from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION

    palette = style.chart_palette
    n_points = max((len(s.values) for s in block.series), default=0)

    for i, series in enumerate(chart.series):
        color = palette[i % len(palette)]
        try:
            if block.chart_type == "pie":  # color each slice, not the series
                for j, point in enumerate(series.points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = palette[j % len(palette)]
            elif block.chart_type in ("line", "scatter"):
                series.format.line.color.rgb = color
                series.format.line.width = Pt(2.25)
            else:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
        except Exception:
            pass

    # Data labels on small bar/pie charts (readable, not cluttered).
    if block.chart_type in ("bar", "pie") and 0 < n_points <= 8 and len(block.series) == 1:
        try:
            plot = chart.plots[0]
            plot.has_data_labels = True
            labels = plot.data_labels
            labels.font.size = Pt(style.chart_axis_pt)
            labels.font.name = style.latin_font
            if block.chart_type == "bar":
                labels.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass

    if block.chart_type == "bar":
        try:
            chart.plots[0].gap_width = 60  # chunkier bars than the 150 default
        except Exception:
            pass

    # Axis fonts + lighter gridlines (category charts; pie has no axes).
    light = RGBColor(0xD9, 0xD9, 0xD9)
    for axis_name in ("category_axis", "value_axis"):
        try:
            axis = getattr(chart, axis_name)
            axis.tick_labels.font.size = Pt(style.chart_axis_pt)
            axis.tick_labels.font.name = style.latin_font
            if axis.has_major_gridlines:
                axis.major_gridlines.format.line.color.rgb = light
        except Exception:
            pass

    chart.has_legend = len(block.series) > 1 or block.chart_type == "pie"
    if chart.has_legend:
        try:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(style.chart_axis_pt)
            chart.legend.font.name = style.latin_font
        except Exception:
            pass


def render_chart(slide, block: ChartBlock, region: Region, style: StyleProfile = ACADEMIC):
    """Render a native, editable PowerPoint chart (not an image), styled by the StyleProfile."""
    left, top, width, height = region
    if block.chart_type == "scatter":
        data = XyChartData()
        all_x: list[float] = []
        all_y: list[float] = []
        for s in block.series:
            series = data.add_series(s.name or "series")
            xs = s.x or list(range(1, len(s.values) + 1))
            for x, y in zip(xs, s.values):
                series.add_data_point(float(x), float(y))
                all_x.append(float(x))
                all_y.append(float(y))
        if block.reference_line and all_x and all_y:  # y=x agreement line spanning the data range
            lo = min(min(all_x), min(all_y))
            hi = max(max(all_x), max(all_y))
            ref = data.add_series("1:1")
            ref.add_data_point(lo, lo)
            ref.add_data_point(hi, hi)
        frame = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, left, top, width, height, data)
        if block.reference_line:
            _style_reference_line(frame.chart, style)
    else:
        n = len(block.categories) if block.categories else max(len(s.values) for s in block.series)
        cats = block.categories or [str(i + 1) for i in range(n)]
        data = CategoryChartData()
        data.categories = cats
        for s in block.series:
            vals = (list(s.values)[:n] + [0.0] * n)[:n]  # pad/truncate to the category count
            data.add_series(s.name or "series", vals)
        ctype = _CATEGORY_CHARTS.get(block.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        frame = slide.shapes.add_chart(ctype, left, top, width, height, data)
    chart = frame.chart
    if block.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = block.title
        try:
            run = chart.chart_title.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(style.table_header_pt)
            run.font.name = style.ea_font
        except Exception:
            pass
    _style_chart(chart, block, style)
    return frame


def _embed_native_formula(slide, latex: str, omml: str, region: Region):
    """Inject an editable OMML equation, wrapped in mc:AlternateContent with a LaTeX-text Fallback so a
    reader that ignores the a14 extension degrades to text (never a blank/corrupt slide)."""
    from xml.sax.saxutils import escape as _esc

    from pptx.oxml import parse_xml

    left, top, width, height = region
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    para = box.text_frame.paragraphs[0]._p
    xml = (
        '<mc:AlternateContent '
        'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        'xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" '
        'xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<mc:Choice Requires="a14"><a14:m><m:oMathPara>{omml}</m:oMathPara></a14:m></mc:Choice>'
        f'<mc:Fallback><a:r><a:t>{_esc(latex)}</a:t></a:r></mc:Fallback>'
        '</mc:AlternateContent>'
    )
    para.append(parse_xml(xml))
    return box


def render_formula(slide, block: FormulaBlock, region: Region, renderer: FormulaRenderer):
    left, top, width, height = region
    # Experimental native-editable OMML (opt-in via the renderer); any failure falls through to image.
    to_omml = getattr(renderer, "to_omml", None)
    if callable(to_omml):
        try:
            omml = to_omml(block.latex)
            if omml:
                return _embed_native_formula(slide, block.latex, omml, region)
        except Exception:
            pass
    image = None
    try:
        image = renderer.to_image(block.latex)
    except Exception:
        image = None
    if image is not None and Path(image).exists():
        return slide.shapes.add_picture(str(image), left, top, width=width)

    # Text fallback: place the LaTeX as editable, centered text.
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = block.latex
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    return box


def render_figure(slide, block: FigureBlock, region: Region, asset_resolver=None, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    # Resolve the asset_id to a rendered image path (Evidence Pool), then fall back to a raw path.
    resolved = asset_resolver.get(block.asset_id) if asset_resolver else None
    candidate = Path(resolved) if resolved else Path(block.asset_id)
    if candidate.is_file():
        caption_h = int(Pt(28)) if block.caption else 0  # reserve a line for the caption
        avail_h = max(1, height - caption_h)
        # Fit the image into (width, avail_h) preserving aspect ratio; center horizontally.
        try:
            from PIL import Image

            with Image.open(str(candidate)) as im:
                iw, ih = im.size
            scale = min(width / iw, avail_h / ih)
            w, h = max(1, int(iw * scale)), max(1, int(ih * scale))
        except Exception:
            w, h = width, avail_h  # defensive (python-pptx already requires Pillow)
        x = left + (width - w) // 2
        y = top + max(0, (avail_h - h) // 2)  # center vertically too (matters in side columns)
        pic = slide.shapes.add_picture(str(candidate), x, y, width=w, height=h)
        if block.caption:
            cap = slide.shapes.add_textbox(left, y + h, width, caption_h)
            cap.text_frame.word_wrap = True
            para = cap.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            add_rich_text(para, block.caption, size=Pt(style.caption_pt), style=style)
        return pic

    # Placeholder when the asset is not resolvable (e.g. the planner described a figure with no asset).
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    text = f"[figure: {block.asset_id}]"
    if block.caption:
        text += f"\n{block.caption}"
    add_rich_text(tf.paragraphs[0], text, size=Pt(14), style=style)
    return box

def _style_reference_line(chart, style: StyleProfile) -> None:
    """Make the last series (the 1:1 line) a thin dashed grey line with no markers; data series
    keep markers only (no connecting line) so the agreement against y=x reads clearly."""
    try:
        from pptx.oxml.ns import qn as _qn

        sers = chart.series
        for i, ser in enumerate(sers):
            is_ref = i == len(sers) - 1
            spPr = ser.format._element.get_or_add_spPr()
            for tag in ("a:ln",):
                for el in spPr.findall(_qn(tag)):
                    spPr.remove(el)
            ln = spPr.makeelement(_qn("a:ln"), {"w": "19050" if is_ref else "0"})
            if is_ref:
                fill = ln.makeelement(_qn("a:solidFill"), {})
                clr = fill.makeelement(_qn("a:srgbClr"), {"val": str(style.muted_rgb)})
                fill.append(clr)
                ln.append(fill)
                dash = ln.makeelement(_qn("a:prstDash"), {"val": "dash"})
                ln.append(dash)
            else:
                ln.append(ln.makeelement(_qn("a:noFill"), {}))
            spPr.append(ln)
    except Exception:
        pass
