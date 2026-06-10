"""Deterministic diagram layout → native, editable PowerPoint shapes.

The LLM emits only a *semantic* `DiagramBlock` (nodes + edges + type, NO coordinates); this module
computes the geometry per `diagram_type` and renders rounded-rectangle nodes + arrow connectors. Because
layout is deterministic, there is no coordinate hallucination and no need for a VLM geometry critic.
"""

from __future__ import annotations

import math
from collections import deque

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from slide_ir import DiagramBlock

from .blocks import add_rich_text
from .style import ACADEMIC, StyleProfile

Region = tuple[int, int, int, int]


def _node(slide, left, top, w, h, label, style: StyleProfile):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(left), int(top), int(max(w, 1)), int(max(h, 1)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = style.node_fill_rgb
    shape.line.color.rgb = style.node_line_rgb
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    # Explicit color: shape-style text defaults to white, invisible on the light node fill.
    add_rich_text(para, label, size=Pt(12), style=style, color=style.text_rgb)
    return shape


def _arrow(slide, x1, y1, x2, y2, style: StyleProfile):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    cxn.line.color.rgb = style.node_line_rgb
    cxn.line.width = Pt(1.5)
    try:  # add an arrowhead (cosmetic) — never let it break rendering
        ln = cxn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    except Exception:
        pass
    return cxn


def _edge_pairs(block: DiagramBlock) -> list[tuple[str, str]]:
    if block.edges:
        return [(e.source, e.target) for e in block.edges]
    nodes = block.nodes
    return [(nodes[i].id, nodes[i + 1].id) for i in range(len(nodes) - 1)]  # sequential fallback


def _flow(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n = len(block.nodes)
    gap = int(width * 0.04)
    box_w = (width - gap * (n - 1)) // max(n, 1)
    box_h = min(int(height * 0.55), int(box_w * 0.7))
    cy = top + height // 2
    boxes = {}
    for i, node in enumerate(block.nodes):
        x = left + i * (box_w + gap)
        y = cy - box_h // 2
        _node(slide, x, y, box_w, box_h, node.label, style)
        boxes[node.id] = (x, y, box_w, box_h)
    for s, t in _edge_pairs(block):
        a, b = boxes.get(s), boxes.get(t)
        if a and b:
            _arrow(slide, a[0] + a[2], a[1] + a[3] // 2, b[0], b[1] + b[3] // 2, style)


def _comparison(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n = len(block.nodes)
    gap = int(width * 0.04)
    box_w = (width - gap * (n - 1)) // max(n, 1)
    for i, node in enumerate(block.nodes):
        _node(slide, left + i * (box_w + gap), top, box_w, int(height * 0.92), node.label, style)


def _cycle(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n = len(block.nodes)
    cx, cy = left + width // 2, top + height // 2
    rx, ry = int(width * 0.33), int(height * 0.33)
    box_w = int(min(width, height) * 0.28)
    box_h = int(box_w * 0.55)
    centers = {}
    for i, node in enumerate(block.nodes):
        ang = -math.pi / 2 + 2 * math.pi * i / n
        x = int(cx + rx * math.cos(ang) - box_w / 2)
        y = int(cy + ry * math.sin(ang) - box_h / 2)
        _node(slide, x, y, box_w, box_h, node.label, style)
        centers[node.id] = (x + box_w // 2, y + box_h // 2)
    pairs = _edge_pairs(block) if block.edges else [
        (block.nodes[i].id, block.nodes[(i + 1) % n].id) for i in range(n)
    ]
    for s, t in pairs:
        if s in centers and t in centers:
            _arrow(slide, *centers[s], *centers[t], style)


def _tree(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    ids = [nd.id for nd in block.nodes]
    labels = {nd.id: nd.label for nd in block.nodes}
    children: dict[str, list[str]] = {i: [] for i in ids}
    indeg = {i: 0 for i in ids}
    for s, t in _edge_pairs(block):
        if s in children and t in indeg:
            children[s].append(t)
            indeg[t] += 1
    roots = [i for i in ids if indeg[i] == 0] or [ids[0]]

    level: dict[int, list[str]] = {}
    seen: set[str] = set()
    q = deque((r, 0) for r in roots)
    while q:
        nid, lv = q.popleft()
        if nid in seen:
            continue
        seen.add(nid)
        level.setdefault(lv, []).append(nid)
        for c in children.get(nid, []):
            q.append((c, lv + 1))
    for i in ids:  # cycles / disconnected -> level 0
        if i not in seen:
            level.setdefault(0, []).append(i)
            seen.add(i)

    nlev = max(level) + 1
    band_h = height // nlev
    box_h = min(int(band_h * 0.6), int(width * 0.12))
    centers = {}
    for lv, nodes in level.items():
        cell = width // max(len(nodes), 1)
        box_w = min(int(cell * 0.85), int(width * 0.3))
        for j, nid in enumerate(nodes):
            x = left + j * cell + (cell - box_w) // 2
            y = top + lv * band_h + (band_h - box_h) // 2
            _node(slide, x, y, box_w, box_h, labels[nid], style)
            centers[nid] = (x + box_w // 2, y, y + box_h)
    for s, t in _edge_pairs(block):
        if s in centers and t in centers:
            _arrow(slide, centers[s][0], centers[s][2], centers[t][0], centers[t][1], style)


def _pyramid(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n = len(block.nodes)
    band = height // max(n, 1)
    for i, node in enumerate(block.nodes):
        factor = (i + 1) / n  # widest at the bottom
        w = int(width * factor)
        _node(slide, left + (width - w) // 2, top + i * band + int(band * 0.1), w, int(band * 0.8), node.label, style)


def _timeline(slide, block, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    n = len(block.nodes)
    cy = top + height // 2
    _arrow(slide, left, cy, left + width, cy, style)  # the timeline axis
    box_w = int(width / max(n, 1) * 0.8)
    box_h = int(height * 0.3)
    for i, node in enumerate(block.nodes):
        cx = left + int(width * (i + 0.5) / n)
        above = i % 2 == 0
        y = cy - box_h - int(height * 0.06) if above else cy + int(height * 0.06)
        _node(slide, cx - box_w // 2, y, box_w, box_h, node.label, style)


_LAYOUTS = {
    "flow": _flow,
    "comparison": _comparison,
    "cycle": _cycle,
    "tree": _tree,
    "pyramid": _pyramid,
    "timeline": _timeline,
}


def render_diagram(slide, block: DiagramBlock, region: Region, style: StyleProfile = ACADEMIC):
    left, top, width, height = region
    inner_top, inner_h = top, height
    if block.title:
        box = slide.shapes.add_textbox(left, top, width, int(Pt(22)))
        para = box.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        add_rich_text(para, block.title, size=Pt(13), bold=True, style=style)
        title_h = int(Pt(26))
        inner_top, inner_h = top + title_h, height - title_h
    _LAYOUTS.get(block.diagram_type, _flow)(slide, block, (left, inner_top, width, inner_h), style)
