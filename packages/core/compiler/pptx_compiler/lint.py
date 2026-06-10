"""Deterministic post-compile geometry lint (AI-free, exact).

Because we own the compiler, layout defects are *computable* from the compiled artifact — no VLM
needed for the basics: text auto-shrunk to the readability floor, figures rendered tiny on
figure-led layouts, content shapes overlapping. Findings are phrased repair-routable
(``slide '<id>': ...``) so the critic retry loop can fix them at the IR level (shorten bullets,
switch layout, split the slide).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from slide_ir import Deck

FONT_FLOOR_PT = 10.0  # blocks.py _fit_font floor — hitting it means the text did not fit at body size
MIN_FIGURE_AREA_FRAC = 0.10  # a figure-led slide whose largest picture is below this is "tiny figure"
MAX_OVERLAP_FRAC = 0.04  # of slide area; content shapes overlapping more than this get flagged
_DECOR_AREA_FRAC = 0.01  # shapes smaller than this are chrome (accent bar, page number) — ignored

_FIGURE_LAYOUTS = {"figure_caption", "figure_left", "big_figure", "figure_grid"}


def _bbox(sh) -> tuple[int, int, int, int]:
    return (sh.left, sh.top, sh.left + sh.width, sh.top + sh.height)


def _overlap(a: tuple, b: tuple) -> int:
    w = min(a[2], b[2]) - max(a[0], b[0])
    h = min(a[3], b[3]) - max(a[1], b[1])
    return w * h if (w > 0 and h > 0) else 0


def lint_compiled_deck(deck: Deck, pptx_path: str | Path) -> list[str]:
    """Lint a compiled deck against its IR. Returns repair-routable findings (empty == clean)."""
    prs = Presentation(str(pptx_path))
    slide_area = int(prs.slide_width) * int(prs.slide_height)
    findings: list[str] = []

    for ir, slide in zip(deck.slides, prs.slides):
        if ir.layout_type.value == "canvas":
            continue  # premium canvas pages own their geometry (guard + finalize passes cover them)
        tag = f"slide '{ir.slide_id}'"
        shapes = [sh for sh in slide.shapes if (sh.width or 0) > 0 and (sh.height or 0) > 0]
        content = [sh for sh in shapes if sh.width * sh.height > _DECOR_AREA_FRAC * slide_area]

        # 1) Text crammed: a substantial text shape whose runs sit at the auto-shrink floor.
        for sh in content:
            if sh.has_text_frame and len(sh.text_frame.text) > 80:
                sizes = [r.font.size for p in sh.text_frame.paragraphs for r in p.runs if r.font.size]
                if sizes and min(sizes) <= Pt(FONT_FLOOR_PT):
                    findings.append(
                        f"{tag}: text crammed (auto-shrunk to the {FONT_FLOOR_PT:g}pt floor) — "
                        "shorten bullets or split this slide"
                    )
                    break

        # 2) Figure-led slide whose figure renders tiny (e.g. extreme aspect ratio in a small cell).
        if ir.layout_type.value in _FIGURE_LAYOUTS:
            pics = [sh for sh in shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            if pics and max(p.width * p.height for p in pics) < MIN_FIGURE_AREA_FRAC * slide_area:
                findings.append(
                    f"{tag}: the figure renders very small — use layout 'big_figure' or reduce co-located blocks"
                )

        # 3) Content shapes overlapping (cannot happen from our templates; guards future layouts).
        boxes = [_bbox(sh) for sh in content]
        flagged = False
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                if _overlap(boxes[i], boxes[j]) > MAX_OVERLAP_FRAC * slide_area:
                    findings.append(f"{tag}: content shapes overlap — simplify blocks or split this slide")
                    flagged = True
                    break
            if flagged:
                break

    return findings
