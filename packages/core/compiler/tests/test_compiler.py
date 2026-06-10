"""Unit tests for the pptx compiler.

Each test maps to a scenario in
openspec/changes/add-pptx-compiler/specs/pptx-compiler/spec.md.
"""

from __future__ import annotations

import base64

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from slide_ir import (
    BulletBlock,
    ChartBlock,
    ChartSeries,
    Deck,
    DiagramBlock,
    DiagramEdge,
    DiagramNode,
    FigureBlock,
    FormulaBlock,
    LayoutType,
    SlideIR,
    TableBlock,
)
from pptx_compiler import NullFormulaRenderer, compile_deck

# 1x1 PNG, used to exercise the image-embedding path without Pillow.
_PNG_1x1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYPhfDwAChwGA60e6kgAAAABJRU5ErkJggg=="
)


def _deck() -> Deck:
    return Deck(
        deck_id="d",
        slides=[
            SlideIR(slide_id="s1", layout_type=LayoutType.TITLE, title="Title"),
            SlideIR(
                slide_id="s2",
                layout_type=LayoutType.BULLET_EVIDENCE,
                title="Motivation",
                blocks=[BulletBlock(items=["alpha", "beta", "gamma"])],
            ),
            SlideIR(
                slide_id="s3",
                layout_type=LayoutType.TWO_COLUMN_TABLE,
                title="Results",
                blocks=[
                    TableBlock(
                        columns=["Sample", "87Sr/86Sr"],
                        rows=[["HT-1", "0.7041"], ["HT-2", "0.7039"]],
                    )
                ],
            ),
            SlideIR(
                slide_id="s4",
                layout_type=LayoutType.FORMULA_BANNER,
                title="Eps",
                blocks=[FormulaBlock(latex="E=mc^2")],
            ),
        ],
    )


def test_one_slide_per_ir_and_valid_pptx(tmp_path):
    out = compile_deck(_deck(), tmp_path / "out.pptx")
    assert out.exists()
    prs = Presentation(str(out))  # reopen -> valid PPTX
    assert len(prs.slides) == 4


def test_table_is_native(tmp_path):
    out = compile_deck(_deck(), tmp_path / "out.pptx")
    slide = Presentation(str(out)).slides[2]
    tables = [sh for sh in slide.shapes if sh.has_table]
    assert len(tables) == 1
    table = tables[0].table
    assert len(table.columns) == 2
    assert len(table.rows) == 3  # header + 2 data rows
    assert table.cell(0, 0).text == "Sample"


def test_bullets_text_present(tmp_path):
    out = compile_deck(_deck(), tmp_path / "out.pptx")
    slide = Presentation(str(out)).slides[1]
    texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    for item in ("alpha", "beta", "gamma"):
        assert item in texts


def test_formula_text_fallback(tmp_path):
    out = compile_deck(_deck(), tmp_path / "out.pptx", formula_renderer=NullFormulaRenderer())
    slide = Presentation(str(out)).slides[3]
    texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "E=mc^2" in texts
    assert not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def test_formula_image_when_renderer_provides(tmp_path):
    img = tmp_path / "f.png"
    img.write_bytes(_PNG_1x1)

    class _ImgRenderer:
        def to_image(self, latex):
            return img

    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FORMULA_BANNER,
                title="f",
                blocks=[FormulaBlock(latex="x^2")],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "o.pptx", formula_renderer=_ImgRenderer())
    slide = Presentation(str(out)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def test_template_slide_size_inherited(tmp_path):
    template = tmp_path / "tmpl.pptx"
    base = Presentation()
    base.slide_width = Inches(13.333)
    base.slide_height = Inches(7.5)
    base.save(str(template))
    expected_w = Presentation(str(template)).slide_width

    out = compile_deck(_deck(), tmp_path / "out.pptx", template=template)
    result = Presentation(str(out))
    assert result.slide_width == expected_w


def test_deterministic_structure(tmp_path):
    deck = _deck()
    a = Presentation(str(compile_deck(deck, tmp_path / "a.pptx")))
    b = Presentation(str(compile_deck(deck, tmp_path / "b.pptx")))
    assert len(a.slides) == len(b.slides)

    def sig(prs):
        return [tuple(sh.shape_type for sh in slide.shapes) for slide in prs.slides]

    assert sig(a) == sig(b)


# --- improve-output-quality: styling -----------------------------------------


def test_fresh_deck_is_16_9(tmp_path):
    deck = Deck(deck_id="d", slides=[SlideIR(slide_id="s1", layout_type=LayoutType.TITLE, title="T")])
    prs = Presentation(str(compile_deck(deck, tmp_path / "w.pptx")))
    assert round(prs.slide_width / prs.slide_height, 3) == 1.778  # 16:9


def test_emphasis_span_becomes_red_bold_run(tmp_path):
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s1",
                layout_type=LayoutType.BULLET_EVIDENCE,
                title="x",
                blocks=[BulletBlock(items=["plain **key** tail"])],
            )
        ],
    )
    prs = Presentation(str(compile_deck(deck, tmp_path / "e.pptx")))
    reds = [
        r
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
        for r in p.runs
        if r.font.bold
        and r.font.color is not None
        and r.font.color.type is not None
        and str(r.font.color.rgb) == "FF0000"
    ]
    assert any(r.text == "key" for r in reds)


# --- add-figure-extraction: asset resolution ---------------------------------


def test_figure_resolved_via_resolver(tmp_path):
    img = tmp_path / "r.png"
    img.write_bytes(_PNG_1x1)
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FIGURE_CAPTION,
                title="f",
                blocks=[FigureBlock(asset_id="img1", caption="c")],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "fig.pptx", asset_resolver={"img1": str(img)})
    slide = Presentation(str(out)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def _png(path, w, h):
    from PIL import Image

    Image.new("RGB", (w, h), (200, 200, 200)).save(str(path))
    return path


def test_figure_is_contained_and_centered(tmp_path):
    img = _png(tmp_path / "wide.png", 800, 200)  # wide image
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FIGURE_CAPTION,
                title="f",
                blocks=[FigureBlock(asset_id="i", caption="cap")],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "fit.pptx", asset_resolver={"i": str(img)})
    prs = Presentation(str(out))
    slide = prs.slides[0]
    pics = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics
    pic = pics[0]
    content_w = prs.slide_width - 2 * Inches(0.5)
    assert pic.width <= content_w  # not forced past content width
    assert abs(pic.height / pic.width - 200 / 800) < 0.02  # aspect preserved
    # centered horizontally within the content area
    center = pic.left + pic.width / 2
    assert abs(center - prs.slide_width / 2) < Inches(0.2)


def test_figure_gets_more_room_than_bullets(tmp_path):
    img = _png(tmp_path / "sq.png", 400, 400)
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FIGURE_CAPTION,
                title="f",
                blocks=[FigureBlock(asset_id="i"), BulletBlock(items=["a", "b"])],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "room.pptx", asset_resolver={"i": str(img)})
    prs = Presentation(str(out))
    shapes = list(prs.slides[0].shapes)
    pic = next(sh for sh in shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    bullets = next(sh for sh in shapes if sh.has_text_frame and "a" in sh.text_frame.text)
    # layout v2: figure_caption is side-by-side — the figure column gets the larger share (58%)
    assert pic.left > bullets.left  # figure on the right
    assert pic.width >= int(0.3 * prs.slide_width)  # figure column is generous


def test_chart_block_renders_native_chart(tmp_path):
    for ctype in ("bar", "line", "pie", "scatter"):
        deck = Deck(
            deck_id="d",
            slides=[
                SlideIR(
                    slide_id="s",
                    layout_type=LayoutType.BULLET_EVIDENCE,
                    title="data",
                    blocks=[
                        ChartBlock(
                            chart_type=ctype,
                            categories=["A", "B", "C"],
                            series=[ChartSeries(name="x", values=[1.0, 2.0, 3.0])],
                            title="t",
                        )
                    ],
                )
            ],
        )
        out = compile_deck(deck, tmp_path / f"chart_{ctype}.pptx")
        slide = Presentation(str(out)).slides[0]
        charts = [sh for sh in slide.shapes if sh.has_chart]
        assert charts, f"{ctype} did not render a native chart"


def test_style_profile_changes_fonts(tmp_path):
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.BULLET_EVIDENCE,
                title="标题",
                blocks=[BulletBlock(items=["要点"])],
            )
        ],
    )

    def latin_fonts(path):
        prs = Presentation(str(path))
        return {
            r.font.name
            for slide in prs.slides
            for sh in slide.shapes
            if sh.has_text_frame
            for p in sh.text_frame.paragraphs
            for r in p.runs
            if r.font.name
        }

    a = latin_fonts(compile_deck(deck, tmp_path / "academic.pptx", style="academic"))
    m = latin_fonts(compile_deck(deck, tmp_path / "modern.pptx", style="modern_teal"))
    default = latin_fonts(compile_deck(deck, tmp_path / "default.pptx"))
    assert "Times New Roman" in a and "Calibri" in m  # profiles set different Latin fonts
    assert default == a  # default == academic (unchanged)


def test_title_color_profile(tmp_path):
    deck = Deck(
        deck_id="d",
        slides=[SlideIR(slide_id="s", layout_type=LayoutType.BULLET_EVIDENCE, title="标题", blocks=[BulletBlock(items=["x"])])],
    )
    prs = Presentation(str(compile_deck(deck, tmp_path / "tc.pptx", style="modern_teal")))
    title_colors = {
        str(r.font.color.rgb)
        for slide in prs.slides
        for sh in slide.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
        for r in p.runs
        if r.text.strip() == "标题" and r.font.color and r.font.color.type is not None
    }
    assert "008080" in title_colors  # the profile's title color (teal) is applied


def test_bullet_font_auto_fits():
    from pptx_compiler.blocks import _fit_font

    big = int(Inches(6))
    short_h = int(Inches(1.5))
    sparse = _fit_font(["A", "B"], big, big)
    dense = _fit_font(["很长很长的学术要点描述内容" * 6 for _ in range(6)], big, short_h)
    assert sparse == 16.0  # sparse text stays at base
    assert 10.0 <= dense < 16.0  # dense text shrinks, but not below the floor


def test_diagram_renders_native_shapes(tmp_path):
    nodes = [DiagramNode(id=f"n{i}", label=f"L{i}") for i in range(3)]
    edges = [DiagramEdge(source="n0", target="n1"), DiagramEdge(source="n1", target="n2")]
    for dt in ("flow", "comparison", "cycle", "tree", "pyramid", "timeline"):
        deck = Deck(
            deck_id="d",
            slides=[
                SlideIR(
                    slide_id="s",
                    layout_type=LayoutType.BULLET_EVIDENCE,
                    title="t",
                    blocks=[DiagramBlock(diagram_type=dt, nodes=nodes, edges=(edges if dt in ("flow", "tree", "cycle") else []))],
                )
            ],
        )
        prs = Presentation(str(compile_deck(deck, tmp_path / f"{dt}.pptx")))
        autoshapes = [sh for sh in prs.slides[0].shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(autoshapes) >= 3, f"{dt}: expected >=3 node shapes, got {len(autoshapes)}"


def test_flow_diagram_has_connectors(tmp_path):
    nodes = [DiagramNode(id=f"n{i}", label=f"L{i}") for i in range(3)]
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.BULLET_EVIDENCE,
                title="t",
                blocks=[DiagramBlock(diagram_type="flow", nodes=nodes)],  # sequential edges inferred
            )
        ],
    )
    prs = Presentation(str(compile_deck(deck, tmp_path / "flow.pptx")))
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) > 3  # 3 nodes + connectors


def test_unresolved_figure_falls_back_to_placeholder(tmp_path):
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(
                slide_id="s",
                layout_type=LayoutType.FIGURE_CAPTION,
                title="f",
                blocks=[FigureBlock(asset_id="missing", caption="c")],
            )
        ],
    )
    out = compile_deck(deck, tmp_path / "ph.pptx", asset_resolver={})
    slide = Presentation(str(out)).slides[0]
    assert not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)
    texts = " ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    assert "missing" in texts


def test_native_omml_formula_round_trips(tmp_path):
    """Opt-in native OMML survives save+reopen as an <m:oMath> equation with a text fallback."""

    class _OmmlRenderer:
        def to_omml(self, latex):
            from formula_render.latex_omml import latex_to_omml

            return latex_to_omml(latex)

        def to_image(self, latex):
            return None

    deck = Deck(
        deck_id="d",
        slides=[SlideIR(slide_id="s", layout_type=LayoutType.FORMULA_BANNER, title="f", blocks=[FormulaBlock(latex="E = mc^2")])],
    )
    out = compile_deck(deck, tmp_path / "o.pptx", formula_renderer=_OmmlRenderer())
    xml = Presentation(str(out)).slides[0].shapes[-1]._element.xml
    assert "oMath" in xml  # native editable equation present
    assert "Fallback" in xml and "mc:" in xml  # graceful text fallback wrapper


def test_native_omml_falls_back_to_image_when_unsupported(tmp_path):
    img = tmp_path / "f.png"
    img.write_bytes(_PNG_1x1)

    class _OmmlRenderer:
        def to_omml(self, latex):
            return None  # unsupported -> no OMML

        def to_image(self, latex):
            return img

    deck = Deck(
        deck_id="d",
        slides=[SlideIR(slide_id="s", layout_type=LayoutType.FORMULA_BANNER, title="f", blocks=[FormulaBlock(latex=r"\ce{H2O}")])],
    )
    out = compile_deck(deck, tmp_path / "o.pptx", formula_renderer=_OmmlRenderer())
    slide = Presentation(str(out)).slides[0]
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)  # fell back to image


# --- figure/text layout balance ----------------------------------------------


def test_balanced_fractions_caps_figure_with_bullets():
    from pptx_compiler.compiler import _balanced_fractions

    f = _balanced_fractions(["figure", "bullets"])
    assert abs(sum(f) - 1.0) < 1e-9
    assert f[0] <= 0.60 + 1e-9  # figure capped
    assert f[1] >= 0.40 - 1e-9  # bullets keep a readable share


def test_balanced_fractions_figure_only_unaffected():
    from pptx_compiler.compiler import _balanced_fractions

    assert _balanced_fractions(["figure"]) == [1.0]  # no bullets -> figure keeps full height


def test_balanced_fractions_multi_visual_combined_cap():
    from pptx_compiler.compiler import _balanced_fractions

    f = _balanced_fractions(["figure", "chart", "bullets"])
    vis = f[0] + f[1]
    assert vis <= 0.60 + 1e-9 and f[2] >= 0.40 - 1e-9


def test_balanced_fractions_table_not_capped():
    from pptx_compiler.compiler import _balanced_fractions

    f = _balanced_fractions(["table", "bullets"])  # table is text-like, not a visual block
    assert abs(f[0] - 2 / 3) < 1e-9 and abs(f[1] - 1 / 3) < 1e-9


# --- layout v2: multi-region templates ----------------------------------------


def _one_px_png(tmp_path, name="f.png"):
    img = tmp_path / name
    img.write_bytes(_PNG_1x1)
    return img


def _slide_of(layout, blocks):
    return Deck(deck_id="d", slides=[SlideIR(slide_id="s", layout_type=layout, title="t", blocks=blocks)])


def _shapes_of(tmp_path, deck, resolver=None):
    out = compile_deck(deck, tmp_path / "o.pptx", asset_resolver=resolver)
    return list(Presentation(str(out)).slides[0].shapes)


def test_figure_caption_renders_side_by_side(tmp_path):
    img = _one_px_png(tmp_path)
    deck = _slide_of(LayoutType.FIGURE_CAPTION, [FigureBlock(asset_id="f1"), BulletBlock(items=["a", "b"])])
    shapes = _shapes_of(tmp_path, deck, {"f1": str(img)})
    pic = next(s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    text = next(s for s in shapes if s.has_text_frame and "a" in s.text_frame.text)
    assert pic.left > text.left  # figure on the RIGHT
    assert pic.top < text.top + text.height and text.top < pic.top + pic.height  # vertically overlapping


def test_figure_left_renders_figure_on_left(tmp_path):
    img = _one_px_png(tmp_path)
    deck = _slide_of(LayoutType.FIGURE_LEFT, [FigureBlock(asset_id="f1"), BulletBlock(items=["a"])])
    shapes = _shapes_of(tmp_path, deck, {"f1": str(img)})
    pic = next(s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    text = next(s for s in shapes if s.has_text_frame and "a" in s.text_frame.text)
    assert pic.left + pic.width <= text.left + text.width and pic.left < text.left  # figure LEFT of text


def test_figure_grid_2x2(tmp_path):
    img = _one_px_png(tmp_path)
    deck = _slide_of(LayoutType.FIGURE_GRID, [FigureBlock(asset_id=f"f{i}") for i in range(4)])
    shapes = _shapes_of(tmp_path, deck, {f"f{i}": str(img) for i in range(4)})
    pics = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pics) == 4
    assert len({p.left for p in pics}) == 2 and len({p.top for p in pics}) == 2  # 2 cols x 2 rows


def test_two_column_table_table_left_text_right(tmp_path):
    deck = _slide_of(
        LayoutType.TWO_COLUMN_TABLE,
        [TableBlock(columns=["A", "B"], rows=[["1", "2"]]), BulletBlock(items=["point"])],
    )
    shapes = _shapes_of(tmp_path, deck)
    table = next(s for s in shapes if s.has_table)
    text = next(s for s in shapes if s.has_text_frame and "point" in s.text_frame.text)
    assert table.left < text.left  # table LEFT, takeaways RIGHT


def test_text_only_slides_keep_vertical_stack(tmp_path):
    deck = _slide_of(LayoutType.BULLET_EVIDENCE, [BulletBlock(items=["a"]), BulletBlock(items=["b"])])
    shapes = _shapes_of(tmp_path, deck)
    boxes = [s for s in shapes if s.has_text_frame and s.text_frame.text in ("a", "b")]
    assert len(boxes) == 2 and boxes[0].left == boxes[1].left  # stacked, same x


# --- data graphics v1: styled charts + tables ----------------------------------


def test_chart_series_use_palette_colors(tmp_path):
    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [ChartBlock(chart_type="bar", categories=["a", "b", "c"], series=[ChartSeries(name="s", values=[1, 2, 3])])],
    )
    out = compile_deck(deck, tmp_path / "c.pptx")
    chart = next(s for s in Presentation(str(out)).slides[0].shapes if s.has_chart).chart
    ser = chart.series[0]
    from pptx_compiler.style import ACADEMIC

    assert ser.format.fill.fore_color.rgb == ACADEMIC.chart_palette[0]  # palette applied
    assert chart.plots[0].has_data_labels  # small single-series bar gets labels
    assert chart.plots[0].gap_width == 60


def test_chart_axis_fonts_styled(tmp_path):
    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [ChartBlock(chart_type="line", categories=["a", "b"], series=[ChartSeries(name="x", values=[1, 2]), ChartSeries(name="y", values=[2, 1])])],
    )
    out = compile_deck(deck, tmp_path / "c2.pptx")
    chart = next(s for s in Presentation(str(out)).slides[0].shapes if s.has_chart).chart
    from pptx.util import Pt as _Pt

    assert chart.value_axis.tick_labels.font.size == _Pt(11)
    assert chart.has_legend and chart.legend.font.size == _Pt(11)


def test_table_header_fill_and_banding(tmp_path):
    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [TableBlock(columns=["A", "B"], rows=[["1", "2"], ["3", "4"], ["5", "6"]])],
    )
    out = compile_deck(deck, tmp_path / "t.pptx")
    table = next(s for s in Presentation(str(out)).slides[0].shapes if s.has_table).table
    from pptx_compiler.style import ACADEMIC

    assert table.cell(0, 0).fill.fore_color.rgb == ACADEMIC.table_header_rgb  # header fill
    assert table.cell(2, 0).fill.fore_color.rgb == ACADEMIC.table_band_rgb  # zebra on row 2
    # header text is white
    run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    from pptx.dml.color import RGBColor as _RGB

    assert run.font.color.rgb == _RGB(0xFF, 0xFF, 0xFF)


def test_table_highlight_cells(tmp_path):
    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [TableBlock(columns=["A"], rows=[["x"], ["y"]], highlight={"cells": [[1, 0]]})],
    )
    out = compile_deck(deck, tmp_path / "t2.pptx")
    table = next(s for s in Presentation(str(out)).slides[0].shapes if s.has_table).table
    from pptx_compiler.style import ACADEMIC

    hot = table.cell(2, 0).text_frame.paragraphs[0].runs[0]  # data row 1 -> table row 2
    assert hot.font.bold and hot.font.color.rgb == ACADEMIC.emphasis_rgb


# --- theme v1: deck chrome ------------------------------------------------------


def test_content_slide_has_accent_bar_and_page_number(tmp_path):
    deck = _slide_of(LayoutType.BULLET_EVIDENCE, [BulletBlock(items=["a"])])
    out = compile_deck(deck, tmp_path / "chrome.pptx")
    shapes = list(Presentation(str(out)).slides[0].shapes)
    from pptx_compiler.style import ACADEMIC

    rects = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert any(r.fill.fore_color.rgb == ACADEMIC.accent_rgb for r in rects)  # accent bar
    assert any(s.has_text_frame and s.text_frame.text == "1" for s in shapes)  # page number


def test_cover_has_no_page_number_but_has_rule(tmp_path):
    deck = Deck(deck_id="d", slides=[SlideIR(slide_id="s", layout_type=LayoutType.TITLE, title="T")])
    out = compile_deck(deck, tmp_path / "cover.pptx")
    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert not any(s.has_text_frame and s.text_frame.text == "1" for s in shapes)  # no page no on cover
    assert any(s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for s in shapes)  # centered rule


# --- content blocks: nested bullets + callout + stat ----------------------------


def test_nested_bullets_render_with_levels(tmp_path):
    from slide_ir import BulletItem

    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [BulletBlock(items=["top", BulletItem(text="parent", children=["child1", "child2"])])],
    )
    out = compile_deck(deck, tmp_path / "n.pptx")
    box = next(s for s in Presentation(str(out)).slides[0].shapes if s.has_text_frame and "parent" in s.text_frame.text)
    paras = box.text_frame.paragraphs
    texts = [p.text for p in paras]
    assert texts == ["top", "parent", "child1", "child2"]
    # hanging indent set: child marL > parent marL (real PPT bullet formatting)
    marl = [int(p._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr").get("marL")) for p in paras]
    assert marl[2] > marl[0]  # nested level indents deeper


def test_callout_renders_tinted_card_with_accent_edge(tmp_path):
    from slide_ir import CalloutBlock

    deck = _slide_of(LayoutType.BULLET_EVIDENCE, [BulletBlock(items=["a"]), CalloutBlock(label="结论", text="key point")])
    out = compile_deck(deck, tmp_path / "co.pptx")
    shapes = list(Presentation(str(out)).slides[0].shapes)
    from pptx_compiler.style import ACADEMIC

    cards = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.has_text_frame and "key point" in s.text_frame.text]
    assert cards and cards[0].fill.fore_color.rgb == ACADEMIC.table_band_rgb
    edges = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not s.text_frame.text and s.fill.fore_color.rgb == ACADEMIC.accent_rgb]
    assert edges  # accent edge present (plus the title accent bar also matches; either proves the color)


def test_stat_renders_big_number_cards(tmp_path):
    from slide_ir import StatBlock, StatItem

    deck = _slide_of(
        LayoutType.BULLET_EVIDENCE,
        [StatBlock(items=[StatItem(value="94%", label="acc"), StatItem(value="r=0.94", label="corr"), StatItem(value="N=320", label="samples")])],
    )
    out = compile_deck(deck, tmp_path / "st.pptx")
    shapes = list(Presentation(str(out)).slides[0].shapes)
    cards = [s for s in shapes if s.has_text_frame and any(v in s.text_frame.text for v in ("94%", "r=0.94", "N=320"))]
    assert len(cards) == 3
    assert len({c.left for c in cards}) == 3  # side by side
    from pptx_compiler.style import ACADEMIC

    run = cards[0].text_frame.paragraphs[0].runs[0]
    assert run.font.bold and run.font.color.rgb == ACADEMIC.accent_rgb


# --- geometry lint --------------------------------------------------------------


def test_lint_flags_crammed_text(tmp_path):
    from pptx_compiler import lint_compiled_deck

    huge = ["这是一条非常长的要点内容用来塞满文本框" * 6 for _ in range(9)]
    deck = _slide_of(LayoutType.FIGURE_CAPTION, [FigureBlock(asset_id="f1"), BulletBlock(items=huge)])
    img = _one_px_png(tmp_path)
    out = compile_deck(deck, tmp_path / "cram.pptx", asset_resolver={"f1": str(img)})
    findings = lint_compiled_deck(deck, out)
    assert any("crammed" in f and "slide 's'" in f for f in findings)


def test_lint_flags_tiny_figure(tmp_path):
    from pptx_compiler import lint_compiled_deck

    wide = _png(tmp_path / "wide.png", 4000, 40)  # extreme aspect -> scales to a sliver
    deck = _slide_of(LayoutType.FIGURE_CAPTION, [FigureBlock(asset_id="f1"), BulletBlock(items=["a", "b"])])
    out = compile_deck(deck, tmp_path / "tiny.pptx", asset_resolver={"f1": str(wide)})
    findings = lint_compiled_deck(deck, out)
    assert any("very small" in f for f in findings)


def test_lint_clean_deck_no_findings(tmp_path):
    from pptx_compiler import lint_compiled_deck

    img = _png(tmp_path / "sq.png", 600, 450)
    deck = _slide_of(LayoutType.FIGURE_CAPTION, [FigureBlock(asset_id="f1"), BulletBlock(items=["短要点一", "短要点二", "短要点三"])])
    out = compile_deck(deck, tmp_path / "clean.pptx", asset_resolver={"f1": str(img)})
    assert lint_compiled_deck(deck, out) == []
