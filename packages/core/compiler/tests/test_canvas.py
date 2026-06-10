"""VisualCanvas (premium tier): guard bans + editable injection. Maps to add-visual-canvas spec."""

from __future__ import annotations

from pptx import Presentation

from slide_ir import CanvasBlock, Deck, LayoutType, SlideIR
from pptx_compiler import canvas_engine_available, compile_deck, validate_canvas_svg

_OK = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720">'
    '<rect x="60" y="140" width="500" height="180" fill="#F5F7FA" stroke="#D0D7E0"/>'
    '<text x="90" y="80" font-size="30" fill="#333333" font-family="SimHei">画布标题</text>'
    '<text x="90" y="200" font-size="18" fill="#8B1A1A">关键数值 r=0.938</text>'
    "</svg>"
)


def test_valid_canvas_passes_guard():
    assert validate_canvas_svg(_OK) == []


def test_guard_bans_unsafe_content():
    bad_script = _OK.replace("</svg>", "<script>alert(1)</script></svg>")
    assert any("script" in f for f in validate_canvas_svg(bad_script))
    bad_foreign = _OK.replace("</svg>", "<foreignObject/></svg>")
    assert any("foreignObject" in f for f in validate_canvas_svg(bad_foreign))
    bad_href = _OK.replace("</svg>", '<use href="http://evil/x"/></svg>')
    assert any("href" in f for f in validate_canvas_svg(bad_href))
    bad_image = _OK.replace("</svg>", '<image href="#x"/></svg>')
    assert any("image" in f for f in validate_canvas_svg(bad_image))
    bad_viewbox = _OK.replace('viewBox="0 0 1280 720"', 'viewBox="0 0 100 100"')
    assert any("viewBox" in f for f in validate_canvas_svg(bad_viewbox))
    assert any("XML" in f for f in validate_canvas_svg("<svg><unclosed</svg>"))


def test_canvas_slide_injected_as_editable_shapes(tmp_path):
    if not canvas_engine_available():  # vendored engine is an app-level dep, not a core one
        import pytest

        pytest.skip("asa_svg2pptx not installed")
    deck = Deck(
        deck_id="d",
        slides=[
            SlideIR(slide_id="n", layout_type=LayoutType.BULLET_EVIDENCE, title="普通页"),
            SlideIR(slide_id="cv", layout_type=LayoutType.CANVAS, title="画布", blocks=[CanvasBlock(svg=_OK)]),
        ],
    )
    out = compile_deck(deck, tmp_path / "c.pptx")
    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    texts = " ".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "画布标题" in texts  # converted to native, editable text
    assert "r=0.938" in texts
    # the normal slide is untouched by the injection
    texts0 = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes if sh.has_text_frame)
    assert "普通页" in texts0


def test_invalid_canvas_falls_back_without_breaking_deck(tmp_path):
    bad = _OK.replace("</svg>", "<script>x</script></svg>")
    deck = Deck(
        deck_id="d",
        slides=[SlideIR(slide_id="cv", layout_type=LayoutType.CANVAS, title="坏画布", blocks=[CanvasBlock(svg=bad)])],
    )
    out = compile_deck(deck, tmp_path / "b.pptx")
    assert len(Presentation(str(out)).slides) == 1  # deck still ships; invalid canvas is not injected
