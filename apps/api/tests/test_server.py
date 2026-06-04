"""Default-app wiring + end-to-end tests.

Each test maps to a scenario in
openspec/changes/add-server-entrypoint/specs/server/spec.md.
"""

from __future__ import annotations

import json

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from asa_agents import FakeLLM
from asa_api.server import build_default_app

# build_default_app uses the two-stage detailed planner, so script: skeleton + one slide per plan.
_SKELETON = json.dumps(
    {
        "slides": [
            {"slide_id": "s1", "layout_type": "title", "title": "Demo", "focus": "intro", "evidence_pages": [], "figure_id": None},
            {"slide_id": "s2", "layout_type": "formula_banner", "title": "Energy", "focus": "formula", "evidence_pages": [], "figure_id": None},
        ]
    }
)
_S1 = json.dumps({"slide_id": "s1", "layout_type": "title", "title": "Demo", "blocks": [], "speaker_notes": "x", "provenance": {"source": "p"}})
_S2 = json.dumps(
    {"slide_id": "s2", "layout_type": "formula_banner", "title": "Energy", "blocks": [{"type": "formula", "latex": "E=mc^2"}], "speaker_notes": "x", "provenance": {"source": "p"}}
)


def test_build_default_app_with_injected_llm(tmp_path):
    app = build_default_app(llm=FakeLLM(_SKELETON), out_dir=tmp_path)  # no provider SDK needed
    assert app is not None


def test_end_to_end_produces_formula_picture(tmp_path):
    client = TestClient(build_default_app(llm=FakeLLM(_SKELETON, _S1, _S2), out_dir=tmp_path))
    job_id = client.post("/jobs", json={"inputs": []}).json()["job_id"]
    client.get(f"/jobs/{job_id}/stream")  # plan -> pause
    client.post(f"/jobs/{job_id}/approve", json={"approved": True})  # resume -> compile
    content = client.get(f"/jobs/{job_id}/download").content

    out = tmp_path / "out.pptx"
    out.write_bytes(content)
    prs = Presentation(str(out))
    pictures = [
        sh for slide in prs.slides for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert pictures  # the formula was rendered to an image and embedded natively


def test_main_is_callable():
    import asa_api.__main__ as entry

    assert callable(entry.main)


def test_load_env_is_best_effort(tmp_path):
    from asa_api.server import load_env

    load_env(path=str(tmp_path / "nonexistent.env"))  # must not raise
